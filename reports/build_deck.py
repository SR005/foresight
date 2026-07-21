"""
Build the Project FORESIGHT executive readout (PPTX) from the model outputs.
Run:  python reports/build_deck.py
"""
from pathlib import Path
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parents[1]
PROC = ROOT / "data" / "processed"
FIG = ROOT / "reports" / "figures"

INK = RGBColor(0x1F, 0x1F, 0x3D)
ACCENT = RGBColor(0x6B, 0x6B, 0xD6)
GREY = RGBColor(0x55, 0x55, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xD1, 0x49, 0x5B)

bt = json.load(open(PROC / "backtest_metrics.json"))
rk = json.load(open(PROC / "risk_summary.json"))
ed = json.load(open(PROC / "eda_insights.json"))
CUR = rk["currency"]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = prs.slide_width, prs.slide_height


def slide():
    return prs.slides.add_slide(BLANK)


def box(s, l, t, w, h):
    tb = s.shapes.add_textbox(l, t, w, h); tb.text_frame.word_wrap = True
    return tb.text_frame


def para(tf, text, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
         bullet=False, space=6, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(space)
    r = p.add_run(); r.text = ("•  " + text) if bullet else text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    r.font.name = "Calibri"
    return p


def band(s, color=INK, h=Inches(1.5)):
    sh = s.shapes.add_shape(1, 0, 0, W, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def title_bar(s, text, sub=None):
    band(s, INK, Inches(1.35))
    tf = box(s, Inches(0.6), Inches(0.28), Inches(12), Inches(1.0))
    para(tf, text, 30, WHITE, True, first=True)
    if sub:
        para(tf, sub, 15, RGBColor(0xC9, 0xC9, 0xEE))


def pic(s, name, l, t, w):
    p = FIG / name
    if p.exists():
        s.shapes.add_picture(str(p), l, t, width=w)


# ---------------------------------------------------------------- 1 Title
s = slide(); band(s, INK, H)
tf = box(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(3))
para(tf, "Project FORESIGHT", 54, WHITE, True, first=True)
para(tf, "Demand & Inventory Intelligence for NorthBay Living", 24,
     RGBColor(0xC9, 0xC9, 0xEE))
para(tf, "Executive readout · Data Science engagement", 16,
     RGBColor(0x9E, 0x9E, 0xCF), space=2)
para(tf, "Forecast · stock-risk early warning · planning dashboard · scoring service",
     14, RGBColor(0x9E, 0x9E, 0xCF))

# ---------------------------------------------------------------- 2 Bottom line
s = slide(); title_bar(s, "The bottom line",
                       "What FORESIGHT found in NorthBay's own data")
# three impact cards
cards = [
    (f"{CUR} {rk['capital_locked_overstock']:,.0f}", "capital locked in overstock",
     f"{rk['quadrant_counts'].get('Markdown / clear',0)} SKUs to clear", ACCENT),
    (f"{CUR} {rk['sales_at_risk_stockout']:,.0f}", "sales at risk from stockouts",
     f"{rk['quadrant_counts'].get('Reorder now',0)} SKUs to reorder now", RED),
    (f"{bt['wape_improvement_pct']:.0f}%", "more accurate than the current guess",
     "forecast beats seasonal-naive baseline", RGBColor(0x3A, 0x96, 0x79)),
]
x = Inches(0.6)
for big, mid, small, col in cards:
    c = s.shapes.add_shape(1, x, Inches(1.9), Inches(3.9), Inches(2.6))
    c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0xF2, 0xF2, 0xFA)
    c.line.color.rgb = col; c.line.width = Pt(2); c.shadow.inherit = False
    tf = c.text_frame; tf.word_wrap = True; tf.margin_top = Inches(0.25)
    para(tf, big, 40, col, True, PP_ALIGN.CENTER, first=True)
    para(tf, mid, 16, INK, True, PP_ALIGN.CENTER, space=2)
    para(tf, small, 12, GREY, False, PP_ALIGN.CENTER)
    x += Inches(4.15)
tf = box(s, Inches(0.6), Inches(4.9), Inches(12.1), Inches(2))
para(tf, "Recommendation: clear the top overstocked SKUs to free working capital, "
         "and raise replenishment on the flagged best-sellers before the peak. "
         "Both lists are prioritised by rupees at stake in the dashboard.",
     16, INK, first=True)

# ---------------------------------------------------------------- 3 What we built
s = slide(); title_bar(s, "What we built",
                       "The four asks in the brief — all delivered")
tf = box(s, Inches(0.7), Inches(1.7), Inches(12), Inches(5))
for t in [
    "A weekly, SKU-level demand forecast that beats a naive baseline honestly.",
    "A stockout / overstock early-warning that says what to reorder and what to clear.",
    "A rupee figure on every risk, so Finance can prioritise by cash impact.",
    "A dashboard the ops team can use unaided, plus a deployed scoring service.",
]:
    para(tf, t, 20, INK, bullet=True, space=12)
para(tf, "All of it re-runs end-to-end from the raw extract with one command — "
         "the client can refresh it next month.", 16, GREY, space=12)

# ---------------------------------------------------------------- 4 Data & quality
s = slide(); title_bar(s, "The data, cleaned and trusted")
tf = box(s, Inches(0.7), Inches(1.7), Inches(6.2), Inches(5))
para(tf, f"{ed['total_units']:,} units across {ed['distinct_skus_all']:,} SKUs, "
         "Dec 2009 – Dec 2011.", 18, INK, True, first=True, space=10)
para(tf, "Cleaning removed, with a logged reason for each:", 16, INK, space=6)
for t in ["19,494 cancellations / returns", "6,166 non-positive qty or price rows",
          "33,666 duplicate lines", "4,676 non-product codes (postage, fees)"]:
    para(tf, t, 15, GREY, bullet=True, space=4)
para(tf, "94% of rows retained — a clean, auditable base.", 16, ACCENT, True, space=10)
para(tf, "Note: inventory positions are modelled (the raw feed is sales-only); "
         "on a live engagement they are replaced by NorthBay's stock table.",
     12, GREY, space=8)
pic(s, "fig_weekly_demand.png", Inches(7.1), Inches(2.2), Inches(5.9))

# ---------------------------------------------------------------- 5 Insights
s = slide(); title_bar(s, "What the demand data says")
tf = box(s, Inches(0.7), Inches(1.7), Inches(6.2), Inches(5))
para(tf, f"1.  The top 20% of SKUs drive {ed['top20pct_volume_share']*100:.0f}% "
         "of all units — focus reorder effort here.", 18, INK, True, first=True, space=14)
para(tf, f"2.  {ed['nov_dec_volume_share']*100:.0f}% of the year's volume lands in "
         "Nov–Dec — reorder must lead the peak by the lead time.", 18, INK, True, space=14)
para(tf, f"3.  ~{ed['dead_stock_skus_last8w']} SKUs sold near-zero in the last 8 "
         "weeks — cash sitting in dead stock, ready to clear.", 18, INK, True, space=14)
pic(s, "fig_top_movers.png", Inches(7.1), Inches(1.9), Inches(5.9))

# ---------------------------------------------------------------- 6 Forecast
s = slide(); title_bar(s, "The forecast — and why you can trust it")
tf = box(s, Inches(0.7), Inches(1.7), Inches(5.6), Inches(5))
para(tf, "We compare every model against a seasonal-naive baseline on a fair, "
         "rolling backtest — the honest test for a forecast.", 16, INK, first=True, space=10)
para(tf, f"Model error (WAPE):  {bt['model']['WAPE']:.2f}", 20, ACCENT, True, space=4)
para(tf, f"Baseline error (WAPE):  {bt['seasonal_naive']['WAPE']:.2f}", 18, GREY, True, space=4)
para(tf, f"→ {bt['wape_improvement_pct']:.0f}% more accurate than the baseline.",
     18, RGBColor(0x3A, 0x96, 0x79), True, space=12)
para(tf, "Lower is better. WAPE = total error ÷ total demand. Validated with "
         "rolling-origin backtesting; no future data ever enters a feature.",
     13, GREY)
pic(s, "fig_forecast_example.png", Inches(6.5), Inches(2.0), Inches(6.5))

# ---------------------------------------------------------------- 7 Risk grid
s = slide(); title_bar(s, "From forecast to decision",
                       "Every SKU placed on a stockout vs overstock grid")
pic(s, "fig_decisioning_grid.png", Inches(0.6), Inches(1.6), Inches(5.4))
tf = box(s, Inches(6.4), Inches(1.8), Inches(6.3), Inches(5))
para(tf, "Reorder now — high stockout risk: replenish before it runs out.",
     17, RED, True, first=True, space=10)
para(tf, "Markdown / clear — high overstock: discount to free capital.",
     17, ACCENT, True, space=10)
para(tf, "Healthy — leave as is.", 17, RGBColor(0x3A, 0x96, 0x79), True, space=10)
para(tf, "Watch / volatile — erratic; review by hand.", 17, RGBColor(0xE0, 0xA4, 0x1A),
     True, space=14)
para(tf, "Every SKU carries a recommended action and the rupees at stake, so the "
         "team can triage hundreds at a glance instead of reading a table.",
     14, GREY)

# ---------------------------------------------------------------- 8 Actions
s = slide(); title_bar(s, "Do this first — prioritised by rupees")
tf = box(s, Inches(0.7), Inches(1.6), Inches(6), Inches(5))
para(tf, "Reorder now (top 5)", 18, RED, True, first=True, space=6)
for r in rk["top_reorder"]:
    para(tf, f"{r['sku_id']} · {r['description'][:30]} — {CUR} {r['value_at_stake']:,.0f}",
         13, INK, bullet=True, space=4)
tf2 = box(s, Inches(6.9), Inches(1.6), Inches(6), Inches(5))
para(tf2, "Markdown / clear (top 5)", 18, ACCENT, True, first=True, space=6)
for r in rk["top_markdown"]:
    para(tf2, f"{r['sku_id']} · {r['description'][:30]} — {CUR} {r['value_at_stake']:,.0f}",
         13, INK, bullet=True, space=4)

# ---------------------------------------------------------------- 9 Honesty
s = slide(); title_bar(s, "Accuracy & limitations — the honest version")
tf = box(s, Inches(0.7), Inches(1.7), Inches(12), Inches(5))
for t in [
    f"The model beats the baseline by {bt['wape_improvement_pct']:.0f}% WAPE — real, "
    "but demand is spiky, so treat forecasts as a guided range, not a promise.",
    "Forecasts come with an 80% interval; the further out, the wider it gets.",
    "Inventory positions in this demo are modelled — plug in NorthBay's real stock "
    "table to make the rupee figures exact.",
    "New / long-tail SKUs have little history and are flagged low-confidence, not "
    "silently forecast.",
    "Everything is reproducible and seeded, so the numbers can be re-created and audited.",
]:
    para(tf, t, 17, INK, bullet=True, space=12)

# ---------------------------------------------------------------- 10 How to use
s = slide(); title_bar(s, "How NorthBay uses it")
tf = box(s, Inches(0.7), Inches(1.7), Inches(12), Inches(5))
para(tf, "Planning dashboard — filter by category/SKU, see forecast vs actual and "
         "the prioritised reorder / markdown list. Built for the ops team, no data "
         "scientist required.", 18, INK, bullet=True, first=True, space=12)
para(tf, "Scoring service (API) — ask for any SKU and get its forecast + risk back; "
         "supports single or batch lookups.", 18, INK, bullet=True, space=12)
para(tf, "Monthly refresh — one command re-runs the whole pipeline on the latest "
         "extract.", 18, INK, bullet=True, space=12)
para(tf, "Deliver it like a consultant. Defend it like a scientist.", 16, ACCENT,
     True, PP_ALIGN.CENTER, space=24)

out = ROOT / "reports" / "executive_readout.pptx"
prs.save(str(out))
print("Saved", out, "with", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
